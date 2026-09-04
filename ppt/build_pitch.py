#!/usr/bin/env python
"""Build the qumode ECD+SNAP chemistry pitch deck (Chinese).

Regenerates ppt/qumode-ecd-snap-chemistry-pitch.pptx from scratch.
Facts sourced from:
  - Dutta et al., JCTC 2025, arXiv:2404.10222 (baseline chemistry paper)
  - Dutta et al., arXiv:2501.11735 (constrained optimization, replicated in paper_result/)
  - He, Zgid, Kemper, Freericks, arXiv:2512.21069 (classical reservoir)
  - Zhang & Zhuang, arXiv:2305.01799 (energy-dependent barren plateau)
  - paper_result/out/*.txt (our replication numbers)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy, os

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
OUT = os.path.join(HERE, "qumode-ecd-snap-chemistry-pitch.pptx")

TEXT = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0x00, 0x00, 0x00)
ACCENT = RGBColor(0x00, 0x00, 0x00)
LINK = RGBColor(0x00, 0x00, 0x00)
RED = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def _set_runs(p, segs, size, default_bold=False):
    """segs: list of (text, bold_or_None, color_or_None)."""
    for text, bold, color in segs:
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = default_bold if bold is None else bold
        r.font.color.rgb = TEXT if color is None else color
        # ensure East-Asian font renders consistently
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {})
            rPr.append(ea)
        ea.set('typeface', 'PingFang SC')


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def para(tf, segs, size=14, first=False, space_after=6, space_before=0,
         align=None, line=None, level=0):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    if isinstance(segs, str):
        segs = [(segs, None, None)]
    _set_runs(p, segs, size)
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = level
    if align:
        p.alignment = align
    if line:
        p.line_spacing = line
    return p


def title_bar(s, text, sub=None):
    tb, tf = textbox(s, Inches(0.6), Inches(0.28), Inches(7.9), Inches(0.62))
    para(tf, [(text, True, TEXT)], size=26, first=True, space_after=0)
    if sub:
        tb2, tf2 = textbox(s, Inches(8.55), Inches(0.52), Inches(4.18), Inches(0.4))
        pp = para(tf2, [(sub, False, GRAY)], size=12.5, first=True, space_after=0)
        pp.alignment = PP_ALIGN.RIGHT
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.98), Inches(12.13), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    ln.line.fill.background()
    ln.shadow.inherit = False


def box(s, x, y, w, h):
    """Plain textbox: no fill, no border."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return tb, tf


# ============================================================ S1 标题
s = slide()
tb, tf = textbox(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(1.9))
para(tf, [("基于 ECD + SNAP 的 Qumode 量子化学", True, TEXT)], size=32, first=True, space_after=4)
para(tf, [("打造", True, TEXT), ("整体优化", True, ACCENT), ("的玻色子基态制备流程（qumode VQE）", True, TEXT)], size=32, space_after=0)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(3.35), Inches(3.1), Pt(1))
bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00); bar.line.fill.background(); bar.shadow.inherit = False

tb, tf = textbox(s, Inches(0.7), Inches(3.75), Inches(12.0), Inches(1.6))
para(tf, [("研究主线：把 qumode VQE 当成一条流水线，从头到尾整体优化", True, TEXT)], size=17, first=True, space_after=8)
para(tf, [("① 选化学基（编码）   ② 参数初始化   ③ ECD+SNAP 混合 ansatz（核心）   ④ 含噪声模拟", False, GRAY)], size=16, space_after=6)
para(tf, [("兜底方案：系统研究 SNAP / ECD 电路的 VQE 优化地形（barren plateau）", False, GRAY)], size=16, space_after=0)

tb, tf = textbox(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.85))
para(tf, [("出发点（目前唯一的 qumode 电子结构工作）：Dutta et al., JCTC 2025 · Simulating electronic structure on bosonic quantum computers · ", False, GRAY),
          ("arXiv:2404.10222", False, LINK)], size=13, first=True, space_after=0)

# ============================================================ S2 基线工作
s = slide()
title_bar(s, "现状：唯一的基线工作做了什么，还缺什么", "Dutta et al., JCTC 2025 (arXiv:2404.10222)")

_, tf = box(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.0))
para(tf, [("已有（按原文）", True, ACCENT)], size=17, first=True, space_after=8)
para(tf, [("• 编码：", True, None), ("H_elec —(Jordan–Wigner)→ H_Q = Σ_μ g_μ P_μ，再把 N_Q 比特的 Pauli 词嵌入 qumode Fock 空间，截断 L = 2^(N_Q)", False, None)], size=14, space_after=7)
para(tf, [("• 映射：", True, None), ("每个 Pauli 词经典拟合成 qumode 电路 —— ECD 路线：15 个酉算符的线性组合、深度 N_d=10；SNAP 路线：单个电路、深度 N_d=16", False, None)], size=14, space_after=7)
para(tf, [("• 测量：", True, None), ("Hadamard test 逐项测 Re⟨ψ|U_μ|ψ⟩，经典求和得 ⟨H⟩", False, None)], size=14, space_after=7)
para(tf, [("• H₂ / STO-3G：", True, None), ("4 自旋轨道 → 4 比特 → 1 个 qumode（L=16）+ 2 个 transmon；制备深度 ECD D=9 vs SNAP D=4；测量酉算符 120 个 vs 14 个，两者都达化学精度", False, None)], size=14, space_after=7)
para(tf, [("• 线性 H₄：", True, None), ("8 比特、185 个 Pauli 词 → 2 个 qumode（各 L=16）+ beam-splitter；SNAP-VQE（D=20）在强关联区打败 UCCSD-VQE", False, None)], size=14, space_after=0)

_, tf = box(s, Inches(6.8), Inches(1.2), Inches(5.93), Inches(5.0))
para(tf, [("缺什么（我们的机会）", True, RED)], size=17, first=True, space_after=8)
para(tf, [("• 分子太少：", True, None), ("只有 H₂ 和线性 H₄（类 Hubbard 的氢链基准），且基组固定为 STO-3G 正则轨道 —— 轨道基这个自由度完全没动", False, None)], size=14, space_after=8)
para(tf, [("• 门族没混合：", True, None), ("ECD+旋转 与 SNAP+位移 各跑各的，从未组合在同一电路里", False, None)], size=14, space_after=8)
para(tf, [("• 初始化没研究：", True, None), ("参数直接随机抽取，没有报告成功率对初值的敏感性，规模变大后是否可训练未知", False, None)], size=14, space_after=8)
para(tf, [("• 全部无噪模拟：", True, None), ("光子损耗等真实噪声的影响原文明确留作 future work", False, None)], size=14, space_after=8)
para(tf, [("• 优化地形留白：", True, None), ("原文原话——“understanding their optimization landscape for VQE … we leave for future development”", False, GRAY)], size=14, space_after=0)

tb, tf = textbox(s, Inches(0.6), Inches(6.45), Inches(12.13), Inches(0.6))
para(tf, [("结论：从编码、初始化、ansatz 到噪声，整条流水线每一环都有明确留白 —— 这正是一个可以整体做透的研究方向。", True, ACCENT)], size=15, first=True, space_after=0)

# ============================================================ S3 提案总览
s = slide()
title_bar(s, "我们的提案：整体优化 qumode VQE 流水线")

tb, tf = textbox(s, Inches(0.6), Inches(1.1), Inches(12.13), Inches(0.72))
para(tf, [("一句话：不只改一个电路，而是对「化学基 → 初始化 → ansatz → 噪声评估」四个环节逐一优化，并用统一基准（氢链等）量化每一环的收益。", False, TEXT)], size=15, first=True, space_after=0, line=1.1)

bx = [Inches(0.6), Inches(3.73), Inches(6.86), Inches(9.99)]
bw = Inches(2.95); by = Inches(1.85); bh = Inches(2.5)
heads = ["① 选化学基（编码层）", "② 参数初始化", "③ ECD+SNAP 混合 ansatz", "④ 含噪声 qutip 模拟"]
bodies = [
    "轨道旋转 U 是免费自由度：换基不改物理，但改变 Pauli 结构与 qumode 电路深度。找让 SNAP/ECD 表示最省的基（正则 / 局域化 / 自然轨道）。",
    "原文参数全随机。系统变大后随机初始化还可靠吗？借鉴逐层生长、几何退火、热启动等策略，给出可复现的初始化配方。",
    "核心方向。ECD 负责“搬运布居”（类 Givens 旋转），SNAP 负责“逐 Fock 态加相位”（类 RZZ）。交替叠加，期望更浅的电路达到同等表达力。",
    "原文全部无噪。用我们现成的 qutip 噪声管线（光子损耗 / Lindblad / T1T2 / 控制误差）回答：哪种 ansatz 在真实 cQED 参数下最稳。",
]
for i in range(4):
    shp, tf = box(s, bx[i], by, bw, bh)
    para(tf, [(heads[i], True, ACCENT)], size=15, first=True, space_after=6)
    para(tf, [(bodies[i], False, None)], size=12.5, space_after=0, line=1.08)

_, tf = box(s, Inches(0.6), Inches(4.6), Inches(12.13), Inches(0.95))
para(tf, [("外层优化器也在射程内：", True, TEXT), ("我们已在约束优化问题上验证 Gibbs 目标函数 f = −ln⟨exp(−ηE)⟩ 显著优于裸能量（ECD-VQE 成功率 34%→84%，n=50）；同一套外层思路可直接迁移到化学 VQE。", False, None)], size=14, first=True, space_after=0, line=1.1)

_, tf = box(s, Inches(0.6), Inches(5.75), Inches(12.13), Inches(1.15))
para(tf, [("兜底方案（最坏情况仍有产出）：", True, RED), ("若混合 ansatz 无明显优势，则系统研究 SNAP / ECD / 混合电路的 VQE 优化地形与 barren plateau —— 这是原文点名的 future work，目前只有 ECD 一侧的部分理论结果（Zhang & Zhuang 2023），本身就是一篇扎实的文章。", False, None)], size=14, first=True, space_after=0, line=1.1)

# ============================================================ S4 方向①
s = slide()
title_bar(s, "方向① 选化学基：让哈密顿量更适合 SNAP/ECD")

_, tf = box(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(2.35))
para(tf, [("问题", True, ACCENT)], size=16, first=True, space_after=6)
para(tf, [("原文直接用正则 HF 分子轨道（H₂ 即 σ_g / σ_u）。但轨道基是一个完全免费的幺正自由度：换基不改变物理谱，却剧烈改变 H_Q 的 Pauli 词数量、系数分布，进而改变 qumode 电路的拟合深度与测量成本。", False, None)], size=14, space_after=8, line=1.12)
para(tf, [("候选基：", True, None), ("正则 HF 轨道 ／ 局域化轨道（Edmiston–Ruedenberg 等）／ 自然轨道 ／ 数值优化的任意旋转 U", False, None)], size=14, space_after=0, line=1.12)

_, tf = box(s, Inches(0.6), Inches(3.75), Inches(6.0), Inches(2.5))
para(tf, [("关键公式", True, ACCENT)], size=15, first=True, space_after=6)
para(tf, [("原子基 → 分子轨道：  φ_p = Σ_μ C_μp · χ_μ", False, None)], size=15, space_after=6)
para(tf, [("轨道旋转（待优化）：  φ̃_p = Σ_q U_qp · φ_q ，  U†U = I", False, None)], size=15, space_after=6)
para(tf, [("选基准则：  min_U  Cost( H_Q(U) )", True, None)], size=15, space_after=6)
para(tf, [("Cost 候选：Pauli 词数 N_H 与系数 1-范数；SNAP/ECD 拟合深度 N_d；Hadamard test 酉算符个数", False, GRAY)], size=13, space_after=0, line=1.1)

_, tf = box(s, Inches(6.8), Inches(1.2), Inches(5.93), Inches(3.1))
para(tf, [("为什么有戏（文献启发）", True, ACCENT)], size=16, first=True, space_after=6)
para(tf, [("• Classical reservoir 方法（He, Zgid, Kemper, Freericks, arXiv:2512.21069）用局域化轨道替代正则 HF 轨道，在氢链、N₂、H₂O 等体系上以显著更浅的电路达到化学精度 —— 证明“换基”能实打实省电路。", False, None)], size=14, space_after=7, line=1.12)
para(tf, [("• 轨道旋转可系统性降低 Hamiltonian 的 1-范数（Koridon et al., PRR 2021），而 1-范数直接控制我们的测量与映射成本。", False, None)], size=14, space_after=7, line=1.12)
para(tf, [("• qumode 特有的一环：Pauli 词 → SNAP/ECD 电路的拟合深度取决于矩阵在 Fock 基下的结构 —— 这是 qubit 文献完全没研究过的“基 ↔ 门族匹配”问题。", False, None)], size=14, space_after=0, line=1.12)

_, tf = box(s, Inches(6.8), Inches(4.45), Inches(5.93), Inches(1.9))
para(tf, [("交付物", True, TEXT)], size=15, first=True, space_after=6)
para(tf, [("以 H₂ / H₄ / H₂O 为基准，扫描轨道基，给出第一张「轨道基 → qumode 资源（深度、测量数、精度）」的定量地图，并回答：哪种基最适合 SNAP，哪种最适合 ECD，哪种最适合混合电路。", False, None)], size=14, space_after=0, line=1.12)

# ============================================================ S5 方向②
s = slide()
title_bar(s, "方向② 参数初始化：随机初值可靠吗？")

_, tf = box(s, Inches(0.6), Inches(1.2), Inches(6.1), Inches(1.75))
para(tf, [("现状", True, ACCENT)], size=16, first=True, space_after=6)
para(tf, [("原文与其配套代码把 ansatz 参数直接随机抽取（如 |β|~U(0,3)，角度~U(0,π)），初态取真空。文中未报告对初值的敏感性；系统规模增大后这种做法是否可训练，完全未知（并与 barren plateau 风险直接挂钩）。", False, None)], size=14, space_after=0, line=1.12)

_, tf = box(s, Inches(0.6), Inches(3.05), Inches(6.1), Inches(1.9))
para(tf, [("我们已有的证据（背包问题，ECD-VQE，BFGS，n=50）", True, ACCENT)], size=14.5, first=True, space_after=6)
para(tf, [("• 裸能量目标下，仅 34% 的随机初值最终把最优比特串顶到峰值；换成随机乘积初态也只有 30% —— 初值不是解药，地形本身多盆地（48% 的初值两种初态都失败）。", False, None)], size=13.5, space_after=5, line=1.1)
para(tf, [("• 说明：拼初始化之前，先要摸清地形；拼初始化之时，需要系统的配方而不是碰运气。", False, None)], size=13.5, space_after=0, line=1.1)

_, tf = box(s, Inches(0.6), Inches(5.1), Inches(6.1), Inches(1.85))
para(tf, [("要研究的配方（借鉴 reservoir 论文附录的成熟做法）", True, TEXT)], size=14.5, first=True, space_after=6)
para(tf, [("• 逐层生长：浅层收敛后加层，新参数加小噪声 [−0.01, 0.01]", False, None)], size=13.5, space_after=4)
para(tf, [("• 几何退火：沿键长正反两遍扫描，相邻几何热启动", False, None)], size=13.5, space_after=4)
para(tf, [("• 常数种子 θ₀ ∈ {π, π/2, …} 与 HF/物理动机初值 vs 纯随机的成功率对比曲线（随体系规模）", False, None)], size=13.5, space_after=0)

pic = s.shapes.add_picture(os.path.join(ROOT, "paper_result/out/init_vacuum_vs_random.png"),
                           Inches(6.95), Inches(1.5), width=Inches(5.75))
tb, tf = textbox(s, Inches(6.95), Inches(4.45), Inches(5.75), Inches(1.05))
para(tf, [("50 组配对试验：真空初态 34% vs 冻结随机乘积初态 30%（能量目标，平均峰值概率 0.31 vs 0.27）。初值微调救不了硬地形 —— 初始化研究必须和地形研究一起做。", False, GRAY)], size=12, first=True, space_after=0, line=1.1)

# ============================================================ S6 方向③（定义与直觉）
s = slide()
title_bar(s, "方向③（核心）ECD+SNAP 混合：定义与直觉")

_, tf = box(s, Inches(0.6), Inches(1.2), Inches(6.35), Inches(3.3))
para(tf, [("两族原生门（cQED 硬件的母语）", True, ACCENT)], size=15, first=True, space_after=7)
para(tf, [("位移门：  D(β) = exp( β·b† − β*·b )", False, None)], size=15, space_after=6)
para(tf, [("ECD 门（回声条件位移）：", True, None)], size=14.5, space_after=3)
para(tf, [("ECD(β) = |1⟩⟨0| ⊗ D(β/2) + |0⟩⟨1| ⊗ D(−β/2)", False, None)], size=15, space_after=3)
para(tf, [("配比特旋转 R(θ,φ) = exp[−i(θ/2)(X·cosφ + Y·sinφ)]，一层 = ECD(β)·[I⊗R(θ,φ)]", False, GRAY)], size=12.5, space_after=8)
para(tf, [("SNAP 门（选数相位）：", True, None)], size=14.5, space_after=3)
para(tf, [("S(θ⃗) = Σ_n exp(iθ_n) |n⟩⟨n| ，  n = 0 … L−1", False, None)], size=15, space_after=3)
para(tf, [("配位移门，一层 = S(θ⃗)·D(α)；对每个 Fock 态独立加相位", False, GRAY)], size=12.5, space_after=0)

_, tf = box(s, Inches(7.15), Inches(1.2), Inches(5.58), Inches(3.3))
para(tf, [("设计直觉：ECD 搬运，SNAP 加相位", True, ACCENT)], size=15, first=True, space_after=7)
para(tf, [("• ECD ≈ Givens 旋转：", True, None), ("条件位移在 Fock 空间搬运布居、注入能量、产生比特–模纠缠（负责“hopping”）。", False, None)], size=13.5, space_after=6, line=1.1)
para(tf, [("• SNAP ≈ RZZ / Z 相位层：", True, None), ("Fock 基对角、逐基矢附相位、不改布居不加能量（负责“干涉”）。", False, None)], size=13.5, space_after=6, line=1.1)
para(tf, [("• 交替「搬运 + 相位」正是 classical reservoir 电路（Givens 层 + ZZ 相位层交替，arXiv:2512.21069）在 qumode 上的对应物 —— 该结构已在氢链等体系用更少门达到化学精度。", False, None)], size=13.5, space_after=6, line=1.1)
para(tf, [("• 注意：Fock 对角 cost 下末尾 SNAP 层梯度恒为零 → 电路应以 ECD 收尾（ECD→SNAP→…→ECD）。", False, GRAY)], size=12.5, space_after=0, line=1.1)

_, tf = box(s, Inches(0.6), Inches(4.7), Inches(12.13), Inches(2.2))
para(tf, [("为什么期望混合更省", True, TEXT)], size=15, first=True, space_after=6)
para(tf, [("• 原文数据已暗示互补：同是 H₂，SNAP 制备深度 D=4、测量 14 项；ECD 深度 D=9、测量 120 项 —— SNAP 表达力强但硬件苛刻，ECD 硬件宽松但电路冗长。", False, None)], size=14, space_after=5, line=1.1)
para(tf, [("• 可训练性红利：ECD+旋转（条件高斯门）已被证明可经典模拟（Lu et al. 2026, arXiv:2603.09233）→ 量子优势必须来自 SNAP 的非高斯性；而 SNAP 不增加电路能量，恰好避开“能量依赖 barren plateau”的能量惩罚 —— 混合电路把表达力和能量这两个旋钮解耦。", False, None)], size=14, space_after=5, line=1.1)
para(tf, [("• 实验设计：ECD-only / SNAP-only / Hybrid 三臂消融，相同参数预算，比较深度、精度、测量成本与梯度方差。", False, None)], size=14, space_after=0, line=1.1)

# ============================================================ S7 方向③（历史与硬件）
s = slide()
title_bar(s, "方向③ 续：为何以前分开，为何现在可行")

_, tf = box(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(3.6))
para(tf, [("历史上分开的三个原因（文献调研结论）", True, ACCENT)], size=16, first=True, space_after=7)
para(tf, [("1. 各自已被证明普适，没有理论压力去混合：", True, None), ("{SNAP, D} 普适（Krastanov et al., PRA 2015）；{ECD, R} 普适（Eickbusch et al., Nat. Phys. 2022）。两套编译器、两个社区。", False, None)], size=14, space_after=7, line=1.12)
para(tf, [("2. 诞生于相反的硬件区间：", True, None), ("SNAP 生于强色散（Heeres 2015：χ/2π ≈ 8.3 MHz，选择性脉冲 Ω ≪ χ）；ECD 生于弱色散高相干腔（Eickbusch 2022：χ/2π ≈ 33 kHz，非选择性脉冲 + 大位移 α₀ ~ 30）。同一块固定 χ 的芯片上，一族自然、另一族别扭。", False, None)], size=14, space_after=7, line=1.12)
para(tf, [("3. 应用场景不同：", True, None), ("SNAP 社区做 Fock 相位控制 / 玻色编码；ECD 社区做相空间控制 / GKP。bosonic VQE 原样继承了这条分界线。", False, None)], size=14, space_after=0, line=1.12)

_, tf = box(s, Inches(6.8), Inches(1.2), Inches(5.93), Inches(3.6))
para(tf, [("关键硬件论点：混合的代价是不对称的", True, ACCENT)], size=15.5, first=True, space_after=7)
para(tf, [("• 在 SNAP 级强 χ 设备（χ/2π ~ 0.5–3 MHz，现代常见）上加 ECD 几乎免费：", True, None), ("ECD 只需 χT < 2π 与足够带宽的比特脉冲，不要求大 α₀（Eickbusch 补充材料 S4 明确指出）。", False, None)], size=14, space_after=7, line=1.12)
para(tf, [("• 反过来不行：", True, None), ("弱 χ 的 ECD 设备做 SNAP 需要几十至上百 μs，超出相干时间。", False, None)], size=14, space_after=7, line=1.12)
para(tf, [("• 真实代价是时长而非切换：", True, None), ("SNAP 一层 ≈ 5–10 倍 ECD 一层的时长 → 每层退相干预算不均衡，这正是消融实验要量化的量。", False, None)], size=14, space_after=7, line=1.12)
para(tf, [("• 趋势佐证：混合振子–比特指令集（arXiv:2407.10381）已把两族门并列为同一 ISA 的指令。", False, GRAY)], size=13, space_after=0, line=1.1)

_, tf = box(s, Inches(0.6), Inches(5.0), Inches(12.13), Inches(1.55))
para(tf, [("结论：", True, TEXT), ("“ECD 与 SNAP 分家”是固定 χ 时代的历史产物，不是物理不兼容。在现代中强 χ 的 cQED 平台上，同一块芯片顺序跑两族门在原理上成立 —— 把它们编进同一个 VQE ansatz 是一步自然、但没人走过的棋。", False, None)], size=15, first=True, space_after=0, line=1.15)

# ============================================================ S8 方向④
s = slide()
title_bar(s, "方向④ 含噪声模拟：原文留白，我们管线现成")

_, tf = box(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(1.35))
para(tf, [("原文原话（结论段）：", True, RED)], size=14, first=True, space_after=4)
para(tf, [("“We leave the exploration of the effects of noise on our qumode approaches to future work.”", False, GRAY)], size=13.5, space_after=0, line=1.1)

_, tf = box(s, Inches(0.6), Inches(2.75), Inches(6.0), Inches(2.0))
para(tf, [("我们已有的噪声管线（qumode_vqe.noise，qutip 实现）", True, ACCENT)], size=14.5, first=True, space_after=6)
para(tf, [("• 光子损耗：论文式 Kraus（κτ ≈ 0.003）或 Lindblad 通道", False, None)], size=13.5, space_after=4)
para(tf, [("• Lindblad 统一模型：热激发 + 腔退相 + transmon T1/T2", False, None)], size=13.5, space_after=4)
para(tf, [("• 相干控制误差：β、θ 的幅度/相位偏差", False, None)], size=13.5, space_after=4)
para(tf, [("• 已在约束优化复现中跑通（含 κ 扫描与直方图对比）", False, None)], size=13.5, space_after=0)

_, tf = box(s, Inches(0.6), Inches(4.95), Inches(6.0), Inches(1.9))
para(tf, [("要回答的问题", True, ACCENT)], size=15, first=True, space_after=6)
para(tf, [("• ECD-only / SNAP-only / Hybrid：谁在真实 cQED 参数下最稳？（SNAP 层更长 → “深度换时长”的净效应必须算）", False, None)], size=13.5, space_after=5, line=1.1)
para(tf, [("• 化学精度（1.6 mHa）在多大 κT 预算内保得住？", False, None)], size=13.5, space_after=5)
para(tf, [("• 噪声 × 地形：CV 版“噪声诱导 barren plateau”尚无人做", False, None)], size=13.5, space_after=0)

_, tf = box(s, Inches(6.8), Inches(1.2), Inches(5.93), Inches(5.65))
para(tf, [("交付物", True, ACCENT)], size=16, first=True, space_after=8)
para(tf, [("1. H₂ / H₄ 在真实硬件参数（χ、T1、T2、κ）下的端到端误差预算表：每个环节吃掉多少化学精度。", False, None)], size=14.5, space_after=8, line=1.15)
para(tf, [("2. 「ansatz × 噪声通道」矩阵图：横轴噪声强度，纵轴能量误差，三种 ansatz 三条线 —— 一张图讲清哪个电路值得上机。", False, None)], size=14.5, space_after=8, line=1.15)
para(tf, [("3. 若混合 ansatz 更浅的优势能折算成更短的总时长，就能把“混合更好”从模拟结论升级成硬件论断 —— 这是通向真实实验合作的敲门砖。", False, None)], size=14.5, space_after=0, line=1.15)

# ============================================================ S9 已有积累
s = slide()
title_bar(s, "已有积累：复现 + 完整代码库（可行性证明）")

_, tf = box(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(3.9))
para(tf, [("已完整复现姊妹论文（同一 Yale 团队的约束优化工作）", True, ACCENT)], size=15.5, first=True, space_after=7)
para(tf, [("• Dutta et al., arXiv:2501.11735 的 Fig.5 / Fig.7：7-bit 背包问题，1 qubit + 2 qumodes（Fock 8×8，N_d=5，40 参数）对照 7-qubit QAOA（p=20，40 参数），BFGS，各 50 组配对试验。", False, None)], size=14, space_after=8, line=1.12)
para(tf, [("• 并在其上做了外层改进 —— Gibbs 目标函数（自适应 η）：", True, None)], size=14, space_after=4)
para(tf, [("f = −ln ⟨ exp(−ηE) ⟩", False, ACCENT)], size=16, space_after=6)
para(tf, [("最优比特串命中率：ECD-VQE 34% → 84%，QAOA 10% → 70%。证明“外层流程优化”本身就有巨大杠杆，同一思路可直接迁移到化学 VQE。", False, None)], size=14, space_after=8, line=1.12)
para(tf, [("• 代码库：ECD / SNAP ansatz、噪声通道、参数布局、测试齐全（github.com/zach102824/qumode，qutip 实现）——四个方向都能立即开工。", False, None)], size=14, space_after=0, line=1.12)

pic = s.shapes.add_picture(os.path.join(ROOT, "paper_result/out/gibbs_vs_energy_bitstring.png"),
                           Inches(6.95), Inches(1.35), width=Inches(5.6))
tb, tf = textbox(s, Inches(6.95), Inches(5.05), Inches(5.6), Inches(0.7))
para(tf, [("50 组配对试验：Gibbs 目标 vs 裸能量，qumode 与 QAOA 两种硬件范式同时受益。", False, GRAY)], size=12, first=True, space_after=0, line=1.1)

tb, tf = textbox(s, Inches(0.6), Inches(6.5), Inches(12.13), Inches(0.55))
para(tf, [("这页的作用：证明我们不是从零开始 —— 基线复现、优化器、噪声模拟全部就位，提案里每个方向都有现成的起跑线。", True, ACCENT)], size=14.5, first=True, space_after=0)

# ============================================================ S10 相关工作 1/2
def related_table(s, rows, col_widths=(2.5, 5.2, 4.4)):
    n = len(rows)
    tbl_shape = s.shapes.add_table(n, 3, Inches(0.6), Inches(1.2), Inches(12.13), Inches(0.4 + 0.62 * (n - 1)))
    tbl = tbl_shape.table
    for j, w in enumerate(col_widths):
        tbl.columns[j].width = Inches(w)
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _set_runs(p, [(cell_text, i == 0, TEXT)], 12 if i else 13)
    return tbl

s = slide()
title_bar(s, "相关工作（1/2）：化学与算法")
rows = [
    ("名称", "一句话成果", "我们怎么用 / 缺口"),
    ("Dutta et al. 2404.10222 · JCTC 2025", "玻色机电子结构 VQE：JWT→Fock，H₂ 与线性 H₄，ECD 与 SNAP 两套 ansatz 分开跑。", "我们的 baseline。混合门、换基、初始化、噪声全是留白。"),
    ("Dutta et al. 2501.11735", "同团队：混合 qubit-qumode 设备解约束优化（背包），ECD-VQE 对照 QAOA。", "我们已完整复现 Fig.5/7 并加 Gibbs 外层改进（34%→84%）。"),
    ("Dutta review 2404.10214", "bosonic 设备量子化学综述：振动、电子、映射与门原语一张地图。", "用来画地图、避免与振动化学重复；电子结构门级路线仍是空档。"),
    ("Dutta QSS-VQE 2509.04727", "激发态 QSS-VQE，门集为 SNAP+位移。", "同门集可对照；不是我们第一目标（基态 + 混合门 + 编码）。"),
    ("He et al. 2512.21069（classical reservoir）", "局域化轨道 + Givens/ZZ 交替层，氢链、N₂、H₂O 更浅电路达化学精度。", "方向①（换基）与方向③（搬运+相位交替结构）的直接启发来源。"),
    ("Huh et al. 1412.8427；Wang et al. 1908.03598", "振动玻色采样；cQED 腔上的振动模拟实验。", "qumode 能做化学的存在性证明；非电子结构基态竞争者。"),
]
related_table(s, rows)

# ============================================================ S11 相关工作 2/2
s = slide()
title_bar(s, "相关工作（2/2）：门、编译与优化地形")
rows = [
    ("名称", "一句话成果", "我们怎么用 / 缺口"),
    ("Heeres et al. 1503.01496", "SNAP 门实验实现与定义（强色散 χ/2π ≈ 8.3 MHz，~1 μs）。", "SNAP 定义出处；SNAP-only / Hybrid 从这里出发。"),
    ("Krastanov et al. 1502.08015", "{SNAP, 位移} 对截断谐振子普适。", "混合门集合法性的理论支柱之一；普适 ≠ 现成化学 ansatz。"),
    ("Eickbusch et al. 2111.06414", "ECD 门实验与定义（弱色散 χ/2π ≈ 33 kHz）；补充材料证明 ECD 在任意 χ 下可行（χT<2π）。", "另一族 native 门；“强 χ 上 ECD 几乎免费”的关键依据。"),
    ("Fösel et al. 2004.14256；Job 2307.11900", "更短 SNAP 序列；SNAP+位移到可执行玻色电路的编译。", "编译与脉冲层可借鉴；简化编码能否缩短编译是开放问题。"),
    ("Zhang & Zhuang 2305.01799", "ECD 电路能量依赖 barren plateau：梯度方差 ~ 1/E^(Mν)（浅层 ν=1，深层 ν=2）。", "兜底论文的理论钩子；只覆盖 ECD + 比特旋转参数 + 态制备 cost。"),
    ("Lu et al. 2603.09233", "条件高斯门集（ECD+旋转）可经典模拟。", "量子优势必须来自 SNAP 非高斯性 → 直接抬高方向③的价值。"),
    ("arXiv 2309.14942", "SNAP+位移单 qudit 可训练性分析。", "多模 SNAP / 混合电路的可训练性仍是空白。"),
]
related_table(s, rows)

# ============================================================ S12 兜底
s = slide()
title_bar(s, "兜底方案：讲清 SNAP/ECD 的优化地形")

_, tf = box(s, Inches(0.6), Inches(1.2), Inches(12.13), Inches(1.15))
para(tf, [("原文原话：", True, RED), ("“An important aspect of the qumode gates we have explored here would be understanding their optimization landscape for VQE, which we leave for future development.”", False, GRAY)], size=14, first=True, space_after=0, line=1.15)

_, tf = box(s, Inches(0.6), Inches(2.55), Inches(6.0), Inches(3.3))
para(tf, [("已知的理论边界", True, ACCENT)], size=16, first=True, space_after=7)
para(tf, [("• Zhang & Zhuang（arXiv:2305.01799）：无限维没有 2-design，用“能量正规化系综”替代 —— ECD 电路梯度方差 ~ 1/E^(Mν)，E 为电路典型能量、M 为模数（浅层 ν=1，深层 ν=2）。", False, None)], size=14, space_after=7, line=1.12)
para(tf, [("• 但只覆盖：ECD 门族、比特旋转参数的梯度、态制备型 cost。", False, None)], size=14, space_after=7)
para(tf, [("• SNAP 侧仅有单 qudit 结果（arXiv:2309.14942），无多模、无混合电路。", False, None)], size=14, space_after=0, line=1.12)

_, tf = box(s, Inches(6.8), Inches(2.55), Inches(5.93), Inches(3.3))
para(tf, [("我们要填的空白（数值 + 可能的定标律）", True, ACCENT)], size=15, first=True, space_after=7)
para(tf, [("• SNAP 相位参数与位移参数的梯度方差如何随深度 / 模数 / 能量定标？", False, None)], size=13.5, space_after=5, line=1.1)
para(tf, [("• 化学 Hamiltonian cost（非态制备型）下地形长什么样？", False, None)], size=13.5, space_after=5)
para(tf, [("• 混合 ECD+SNAP：SNAP 不加能量 —— 能否在同等表达力下避开能量惩罚？", False, None)], size=13.5, space_after=5, line=1.1)
para(tf, [("• Gibbs 目标的 η 是否加宽可训练区间？（完全无人做）", False, None)], size=13.5, space_after=5)
para(tf, [("• 加噪声后：CV 版噪声诱导 BP。", False, None)], size=13.5, space_after=0)

tb, tf = textbox(s, Inches(0.6), Inches(6.15), Inches(12.13), Inches(0.9))
para(tf, [("即使方向③没有胜出的 ansatz，「SNAP / ECD / 混合电路的 VQE 优化地形」也是原文点名、理论只做了一角的问题 —— 用我们现成的数值管线系统扫一遍，本身就是一篇有引用价值的文章。", True, ACCENT)], size=15, first=True, space_after=0, line=1.2)

prs.save(OUT)
print("Saved:", OUT, "slides:", len(prs.slides._sldIdLst))
